import json
from datetime import datetime
from pathlib import Path
from uuid import uuid4

from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph as DocxParagraph
from openpyxl import load_workbook
from pptx import Presentation

from app.config.settings import settings
from app.db.database import get_db
from app.utils.chunk import TextPage, chunk_pages


class DocumentService:
    SUPPORTED_TYPES = {"pdf", "md", "txt", "docx", "xlsx", "pptx"}

    def __init__(self) -> None:
        self.upload_dir = Path(settings.data_dir) / "uploads"
        self.chunk_dir = Path(settings.data_dir) / "chunks"
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        self.chunk_dir.mkdir(parents=True, exist_ok=True)

    def _extract_text(self, file_path: Path, file_type: str) -> str:
        return "\n\n".join(page.content for page in self._extract_pages(file_path, file_type))

    def _extract_pages(self, file_path: Path, file_type: str) -> list[TextPage]:
        if file_type in {"md", "txt"}:
            return [TextPage(content=file_path.read_text(encoding="utf-8", errors="ignore"))]
        if file_type == "pdf":
            reader = PdfReader(str(file_path))
            return [
                TextPage(content=page.extract_text() or "", page_no=idx)
                for idx, page in enumerate(reader.pages, start=1)
            ]
        if file_type == "docx":
            return self._extract_docx_pages(file_path)
        if file_type == "xlsx":
            return self._extract_xlsx_pages(file_path)
        if file_type == "pptx":
            return self._extract_pptx_pages(file_path)
        raise ValueError(f"unsupported_file_type: {file_type}")

    def _extract_docx_pages(self, file_path: Path) -> list[TextPage]:
        """从 .docx 文件中按段落+表格提取文本，每 2000 字符切分为一页"""
        doc = DocxDocument(str(file_path))
        pages: list[TextPage] = []
        current_parts: list[str] = []
        current_page = 1

        for element in _iter_docx_elements(doc):
            if isinstance(element, DocxParagraph):
                text = element.text.strip()
                if not text:
                    continue
                # 标题样式转 Markdown 标题前缀，让下游 section 检测能识别
                style_name = element.style.name or ""
                if style_name.startswith("Heading"):
                    level_str = style_name.replace("Heading ", "")
                    try:
                        text = f"{'#' * int(level_str)} {text}"
                    except ValueError:
                        pass
                current_parts.append(text)
            elif isinstance(element, DocxTable):
                current_parts.append(_format_docx_table(element))

            if sum(len(p) for p in current_parts) > 2000:
                pages.append(TextPage(content="\n\n".join(current_parts), page_no=current_page))
                current_page += 1
                current_parts = []

        if current_parts:
            pages.append(TextPage(content="\n\n".join(current_parts), page_no=current_page))
        return pages

    def _extract_xlsx_pages(self, file_path: Path) -> list[TextPage]:
        """从 .xlsx 文件中提取每个工作表的内容"""
        wb = load_workbook(str(file_path), read_only=True, data_only=True)
        pages: list[TextPage] = []
        for sheet_idx, sheet in enumerate(wb.worksheets, start=1):
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                content = f"## {sheet.title}\n\n" + "\n".join(rows)
                pages.append(TextPage(content=content, page_no=sheet_idx))
        wb.close()
        return pages or [TextPage(content="(空工作簿)")]

    def _extract_pptx_pages(self, file_path: Path) -> list[TextPage]:
        """从 .pptx 文件中提取每页幻灯片的文本"""
        prs = Presentation(str(file_path))
        pages: list[TextPage] = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        texts.append(" | ".join(cells))
            if texts:
                pages.append(TextPage(content="\n".join(texts), page_no=slide_idx))
        return pages or [TextPage(content="(空演示文稿)")]

    def fetch_url_and_process(self, url: str) -> dict:
        """抓取网页内容并处理为文档"""
        import requests
        from bs4 import BeautifulSoup

        resp = requests.get(url, timeout=15, headers={"User-Agent": "RaceAgent/1.0"})
        resp.raise_for_status()

        soup = BeautifulSoup(resp.text, "html.parser")
        # 移除 script 和 style 标签
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        if not text or len(text) < 10:
            raise ValueError("网页内容为空或无法提取文本")

        title = soup.title.string.strip() if soup.title and soup.title.string else url
        file_name = f"{title[:60]}.md"
        content = f"# {title}\n\n来源: {url}\n\n{text}"
        return self.upload_and_process(file_name, content.encode("utf-8"))

    def upload_and_process(self, file_name: str, content: bytes, project_id: str = "") -> dict:
        suffix = Path(file_name).suffix.lower().lstrip(".")
        if suffix not in self.SUPPORTED_TYPES:
            raise ValueError("only pdf/md/txt/docx/xlsx/pptx are supported")

        # 文件内容魔数校验
        if suffix == "pdf":
            if not content or not content[:5].startswith(b"%PDF"):
                raise ValueError("文件内容与扩展名不匹配：PDF 文件必须以 %PDF 开头")
        if suffix in {"docx", "xlsx", "pptx"}:
            if not content or content[:2] != b"PK":
                raise ValueError(f"文件内容与扩展名不匹配：{suffix.upper()} 文件必须是有效的 ZIP 格式")

        document_id = uuid4().hex
        stored_name = f"{document_id}_{file_name}"
        file_path = self.upload_dir / stored_name
        file_path.write_bytes(content)

        created_at = datetime.utcnow().isoformat()
        with get_db() as conn:
            conn.execute(
                """
                INSERT INTO documents(
                    id, file_name, file_path, file_type, parse_status, chunk_count,
                    tags, summary, project_id, created_at
                )
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                """,
                (
                    document_id,
                    file_name,
                    str(file_path),
                    suffix,
                    "PROCESSING",
                    0,
                    "",
                    "",
                    project_id or "",
                    created_at,
                ),
            )

        try:
            pages = self._extract_pages(file_path, suffix)
            chunks = chunk_pages(pages, chunk_size=800, chunk_overlap=100)
            chunk_payload = []
            for chunk in chunks:
                chunk_payload.append(
                    {
                        "chunk_id": f"{document_id}_chunk_{chunk.chunk_index:04d}",
                        "document_id": document_id,
                        "source_file": file_name,
                        "file_type": suffix,
                        "chunk_index": chunk.chunk_index,
                        "total_chunks": len(chunks),
                        "content": chunk.content,
                        "page_no": chunk.page_no,
                        "section": chunk.section,
                        "char_start": chunk.char_start,
                        "char_end": chunk.char_end,
                    }
                )
            (self.chunk_dir / f"{document_id}.json").write_text(
                json.dumps(chunk_payload, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
            status = "READY"
            chunk_count = len(chunk_payload)
        except Exception:
            status = "FAILED"
            chunk_count = 0
            raise
        finally:
            with get_db() as conn:
                conn.execute(
                    "UPDATE documents SET parse_status=?, chunk_count=? WHERE id=?",
                    (status, chunk_count, document_id),
                )

        return {
            "document_id": document_id,
            "file_name": file_name,
            "file_type": suffix,
            "parse_status": status,
            "chunk_count": chunk_count,
        }


def _iter_docx_elements(doc):
    """按文档顺序交错产出段落和表格"""
    for element in doc.element.body:
        tag = element.tag
        if tag.endswith("}p"):
            yield DocxParagraph(element, doc)
        elif tag.endswith("}tbl"):
            yield DocxTable(element, doc)


def _format_docx_table(table) -> str:
    """将 docx 表格转为 Markdown 格式"""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(cells)
    if not rows:
        return ""
    # 列数过多时用管道分隔格式，避免 Markdown 表格过宽
    if len(rows[0]) > 5:
        return "\n".join(" | ".join(row) for row in rows)
    result = []
    result.append("| " + " | ".join(rows[0]) + " |")
    result.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
    for row in rows[1:]:
        escaped = [c.replace("|", "\\|") for c in row]
        result.append("| " + " | ".join(escaped) + " |")
    return "\n".join(result)
