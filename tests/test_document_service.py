import sys
import tempfile
from pathlib import Path

import pytest

PROJECT_ROOT = Path(__file__).resolve().parents[1]
BACKEND_ROOT = PROJECT_ROOT / "backend"
sys.path.insert(0, str(BACKEND_ROOT))

try:
    from docx import Document as DocxDocument  # noqa: E402
    from openpyxl import Workbook  # noqa: E402
    from pptx import Presentation  # noqa: E402
    HAS_DOCX_DEPS = True
except ImportError:
    HAS_DOCX_DEPS = False

pytestmark = pytest.mark.skipif(not HAS_DOCX_DEPS, reason="requires python-docx, openpyxl, python-pptx")


@pytest.fixture
def tmp_dir():
    with tempfile.TemporaryDirectory() as d:
        yield Path(d)


def _create_test_docx(path: Path):
    doc = DocxDocument()
    doc.add_heading("测试标题", level=1)
    doc.add_paragraph("这是第一段内容，用于测试文档解析功能。")
    doc.add_paragraph("这是第二段内容，包含更多的文本数据。")
    doc.add_heading("二级标题", level=2)
    doc.add_paragraph("子章节内容。")
    table = doc.add_table(rows=2, cols=3)
    table.cell(0, 0).text = "姓名"
    table.cell(0, 1).text = "年龄"
    table.cell(0, 2).text = "城市"
    table.cell(1, 0).text = "张三"
    table.cell(1, 1).text = "25"
    table.cell(1, 2).text = "北京"
    doc.save(str(path))


def _create_test_xlsx(path: Path):
    wb = Workbook()
    ws1 = wb.active
    ws1.title = "成绩表"
    ws1.append(["姓名", "语文", "数学", "英语"])
    ws1.append(["张三", 90, 85, 92])
    ws1.append(["李四", 88, 95, 78])
    ws2 = wb.create_sheet("汇总")
    ws2.append(["科目", "平均分"])
    ws2.append(["语文", 89])
    ws2.append(["数学", 90])
    wb.save(str(path))


def _create_test_pptx(path: Path):
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "项目介绍"
    slide1.placeholders[1].text = "这是第一页幻灯片的内容"
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "技术方案"
    slide2.placeholders[1].text = "系统架构设计说明\n功能模块划分"
    prs.save(str(path))


def test_extract_docx_pages(tmp_dir):
    from app.service.document_service import DocumentService

    docx_path = tmp_dir / "test.docx"
    _create_test_docx(docx_path)

    svc = DocumentService()
    pages = svc._extract_pages(docx_path, "docx")

    assert len(pages) >= 1
    content = "\n".join(p.content for p in pages)
    assert "测试标题" in content
    assert "第一段内容" in content
    assert "姓名" in content
    assert "张三" in content


def test_extract_docx_headings(tmp_dir):
    from app.service.document_service import DocumentService

    docx_path = tmp_dir / "test.docx"
    _create_test_docx(docx_path)

    svc = DocumentService()
    pages = svc._extract_pages(docx_path, "docx")
    content = "\n".join(p.content for p in pages)

    assert "# 测试标题" in content
    assert "## 二级标题" in content


def test_extract_xlsx_pages(tmp_dir):
    from app.service.document_service import DocumentService

    xlsx_path = tmp_dir / "test.xlsx"
    _create_test_xlsx(xlsx_path)

    svc = DocumentService()
    pages = svc._extract_pages(xlsx_path, "xlsx")

    assert len(pages) == 2
    content = "\n".join(p.content for p in pages)
    assert "成绩表" in content
    assert "张三" in content
    assert "90" in content
    assert "汇总" in content


def test_extract_pptx_pages(tmp_dir):
    from app.service.document_service import DocumentService

    pptx_path = tmp_dir / "test.pptx"
    _create_test_pptx(pptx_path)

    svc = DocumentService()
    pages = svc._extract_pages(pptx_path, "pptx")

    assert len(pages) == 2
    content = "\n".join(p.content for p in pages)
    assert "项目介绍" in content
    assert "技术方案" in content


def test_extract_unsupported_type(tmp_dir):
    from app.service.document_service import DocumentService

    svc = DocumentService()
    with pytest.raises(ValueError, match="unsupported_file_type"):
        svc._extract_pages(tmp_dir / "test.xyz", "xyz")


def test_magic_number_check_pdf(tmp_dir):
    from app.service.document_service import DocumentService

    svc = DocumentService()
    svc.upload_dir = tmp_dir
    with pytest.raises(ValueError, match="PDF"):
        svc.upload_and_process("fake.pdf", b"not a pdf file")


def test_magic_number_check_docx(tmp_dir):
    from app.service.document_service import DocumentService

    svc = DocumentService()
    svc.upload_dir = tmp_dir
    with pytest.raises(ValueError, match="DOCX"):
        svc.upload_and_process("fake.docx", b"not a zip file")


def test_magic_number_check_xlsx(tmp_dir):
    from app.service.document_service import DocumentService

    svc = DocumentService()
    svc.upload_dir = tmp_dir
    with pytest.raises(ValueError, match="XLSX"):
        svc.upload_and_process("fake.xlsx", b"not a zip file")


def test_magic_number_check_pptx(tmp_dir):
    from app.service.document_service import DocumentService

    svc = DocumentService()
    svc.upload_dir = tmp_dir
    with pytest.raises(ValueError, match="PPTX"):
        svc.upload_and_process("fake.pptx", b"not a zip file")


def test_extract_xlsx_empty_workbook(tmp_dir):
    from app.service.document_service import DocumentService

    wb = Workbook()
    xlsx_path = tmp_dir / "empty.xlsx"
    wb.save(str(xlsx_path))

    svc = DocumentService()
    pages = svc._extract_pages(xlsx_path, "xlsx")
    assert len(pages) >= 1
    assert "空工作簿" in pages[0].content


def test_extract_pptx_empty(tmp_dir):
    from app.service.document_service import DocumentService

    prs = Presentation()
    pptx_path = tmp_dir / "empty.pptx"
    prs.save(str(pptx_path))

    svc = DocumentService()
    pages = svc._extract_pages(pptx_path, "pptx")
    assert len(pages) >= 1
    assert "空演示文稿" in pages[0].content
